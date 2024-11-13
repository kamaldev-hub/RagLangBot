import os
from typing import Dict, Any, List
import filetype
from PIL import Image
import pandas as pd
import openpyxl
from pptx import Presentation
import docx
import fitz  # PyMuPDF
import pdfplumber
import csv
import json
import xml.etree.ElementTree as ET
import base64
from io import BytesIO
import logging
import re
from nltk.tokenize import sent_tokenize

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Universal document processor with enhanced RAG capabilities."""

    def __init__(self):
        self.min_chunk_size = 500
        self.max_chunk_size = 2000
        self.chunk_overlap = 200
        self.min_list_item_size = 100
        self.max_list_items_per_chunk = 5

    def process_file(self, file_path: str) -> Dict[str, Any]:
        try:
            kind = filetype.guess(file_path)
            mime_type = kind.mime if kind else 'application/octet-stream'
            file_ext = os.path.splitext(file_path)[1].lower()

            result = {
                'mime_type': mime_type,
                'extension': file_ext,
                'content': None,
                'text': None,
                'metadata': {},
                'error': None,
                'chunks': [],
                'structure': {}
            }

            if mime_type.startswith('image/'):
                result.update(self._process_image(file_path))
            elif file_ext == '.pdf':
                result.update(self._process_pdf_advanced(file_path))
            elif file_ext in ['.xlsx', '.xls']:
                result.update(self._process_excel(file_path))
            elif file_ext in ['.docx', '.doc']:
                result.update(self._process_word(file_path))
            elif file_ext in ['.pptx', '.ppt']:
                result.update(self._process_powerpoint(file_path))
            elif file_ext in ['.txt', '.csv', '.json', '.xml']:
                result.update(self._process_text(file_path, mime_type))
            elif file_ext in ['.py', '.js', '.java', '.cpp', '.cs', '.php', '.rb', '.swift']:
                result.update(self._process_code(file_path))

            if result.get('content'):
                result['structure'] = self._analyze_document_structure(result['content'])
                result['chunks'] = self._create_smart_chunks(result['content'], result['structure'])

            return result

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            return {
                'error': str(e),
                'mime_type': mime_type if 'mime_type' in locals() else None
            }

    def _analyze_document_structure(self, text: str) -> Dict[str, Any]:
        try:
            paragraphs = text.split('\n\n')
            return {
                'has_numbered_lists': bool(re.search(r'^\s*\d+\.', text, re.MULTILINE)),
                'has_bullet_points': bool(re.search(r'^\s*[\•\-\*]', text, re.MULTILINE)),
                'has_headers': bool(re.search(r'^#+\s', text, re.MULTILINE)) or bool(
                    re.search(r'^[A-Z][^.!?]*[:.]', text, re.MULTILINE)),
                'avg_paragraph_length': sum(len(p) for p in paragraphs) / len(paragraphs) if paragraphs else 0,
                'total_length': len(text),
                'has_tables': '[table]' in text.lower() or '|' in text,
                'has_code_blocks': bool(re.search(r'```.*?```', text, re.DOTALL))
            }
        except Exception as e:
            logger.error(f"Error analyzing document structure: {str(e)}")
            return {}

    def _create_smart_chunks(self, text: str, structure: Dict[str, Any]) -> List[str]:
        try:
            if structure.get('has_numbered_lists'):
                return self._chunk_numbered_list(text)
            elif structure.get('has_headers'):
                return self._chunk_with_headers(text)
            else:
                return self._chunk_regular_text(text)
        except Exception as e:
            logger.error(f"Error creating chunks: {str(e)}")
            return [text]

    def _chunk_numbered_list(self, text: str) -> List[str]:
        chunks = []
        current_items = []
        current_length = 0

        items = re.split(r'(?=^\s*\d+\.)', text, flags=re.MULTILINE)

        for item in items:
            item = item.strip()
            if not item:
                continue

            item_length = len(item)

            if current_length + item_length > self.max_chunk_size or len(
                    current_items) >= self.max_list_items_per_chunk:
                if current_items:
                    chunks.append('\n'.join(current_items))
                current_items = [item]
                current_length = item_length
            else:
                current_items.append(item)
                current_length += item_length

        if current_items:
            chunks.append('\n'.join(current_items))

        return chunks

    def _chunk_with_headers(self, text: str) -> List[str]:
        chunks = []
        lines = text.split('\n')
        current_chunk = []
        current_length = 0
        current_header = None

        for line in lines:
            is_header = bool(re.match(r'^#+\s|^[A-Z][^.!?]*[:.]', line))

            if is_header:
                if current_chunk:
                    chunks.append('\n'.join(current_chunk))
                current_chunk = [line]
                current_length = len(line)
                current_header = line
            else:
                if current_length + len(line) > self.max_chunk_size:
                    chunks.append('\n'.join(current_chunk))
                    current_chunk = [current_header] if current_header else []
                    current_chunk.append(line)
                    current_length = sum(len(l) for l in current_chunk)
                else:
                    current_chunk.append(line)
                    current_length += len(line)

        if current_chunk:
            chunks.append('\n'.join(current_chunk))

        return chunks

    def _chunk_regular_text(self, text: str) -> List[str]:
        chunks = []
        sentences = sent_tokenize(text)
        current_chunk = []
        current_length = 0

        for sentence in sentences:
            sentence_length = len(sentence)

            if current_length + sentence_length > self.max_chunk_size:
                if current_chunk:
                    chunks.append(' '.join(current_chunk))
                if chunks:
                    overlap = ' '.join(current_chunk[-2:])
                    current_chunk = [overlap, sentence]
                    current_length = len(overlap) + sentence_length
                else:
                    current_chunk = [sentence]
                    current_length = sentence_length
            else:
                current_chunk.append(sentence)
                current_length += sentence_length

        if current_chunk:
            chunks.append(' '.join(current_chunk))

        return chunks

    def _process_image(self, file_path: str) -> Dict[str, Any]:
        try:
            with Image.open(file_path) as img:
                buffered = BytesIO()
                img.save(buffered, format=img.format)
                img_str = base64.b64encode(buffered.getvalue()).decode()

                return {
                    'content': img_str,
                    'text': "OCR is currently disabled",
                    'metadata': {
                        'format': img.format,
                        'size': img.size,
                        'mode': img.mode,
                        'ocr_enabled': False
                    }
                }
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {str(e)}")
            return {'error': str(e)}

    def _process_pdf_advanced(self, file_path: str) -> Dict[str, Any]:
        try:
            doc = fitz.open(file_path)
            text_content = []
            images = []
            headers = []
            tables = []

            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    table_data = page.extract_tables()
                    if table_data:
                        tables.extend(table_data)

            for page in doc:
                text = page.get_text()
                if text.strip():
                    text_content.append(text)

                blocks = page.get_text("dict")["blocks"]
                for block in blocks:
                    if "lines" in block:
                        for line in block["lines"]:
                            for span in line["spans"]:
                                if span["size"] > 12:
                                    headers.append(span["text"])

                for img in page.get_images():
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_str = base64.b64encode(image_bytes).decode()
                        images.append(image_str)
                    except Exception as e:
                        logger.warning(f"Could not extract image: {str(e)}")

            full_content = '\n'.join(text_content)

            if tables:
                table_text = "\n\nTables found in document:\n"
                for i, table in enumerate(tables, 1):
                    table_text += f"\nTable {i}:\n"
                    for row in table:
                        table_text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
                full_content += table_text

            return {
                'content': full_content,
                'tables': tables,
                'images': images,
                'headers': headers,
                'metadata': {
                    'pages': len(doc),
                    'title': doc.metadata.get('title', ''),
                    'author': doc.metadata.get('author', ''),
                    'has_tables': bool(tables),
                    'has_images': bool(images),
                    'extraction_method': 'PyMuPDF+pdfplumber'
                }
            }
        except Exception as e:
            logger.error(f"Error in advanced PDF processing: {str(e)}")
            return {'error': str(e)}

    def _process_excel(self, file_path: str) -> Dict[str, Any]:
        try:
            df_dict = pd.read_excel(file_path, sheet_name=None)
            text_content = []

            for sheet_name, df in df_dict.items():
                text_content.append(f"Sheet: {sheet_name}")
                text_content.append(df.to_string())

            return {
                'content': '\n'.join(text_content),
                'metadata': {
                    'sheets': list(df_dict.keys()),
                    'total_sheets': len(df_dict)
                }
            }
        except Exception as e:
            logger.error(f"Error processing Excel {file_path}: {str(e)}")
            return {'error': str(e)}

    def _process_word(self, file_path: str) -> Dict[str, Any]:
        try:
            doc = docx.Document(file_path)
            text_content = []

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # Preserve formatting hints
                    if paragraph.style.name.startswith('Heading'):
                        text_content.append(f"### {paragraph.text} ###")
                    else:
                        text_content.append(paragraph.text)

            return {
                'content': '\n'.join(text_content),
                'metadata': {
                    'paragraphs': len(doc.paragraphs)
                }
            }
        except Exception as e:
            logger.error(f"Error processing Word document {file_path}: {str(e)}")
            return {'error': str(e)}

    def _process_powerpoint(self, file_path: str) -> Dict[str, Any]:
        try:
            prs = Presentation(file_path)
            text_content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                slide_content.append(f"Slide {slide_num}:")

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            slide_content.append(shape.text)

                if slide_content:
                    text_content.extend(slide_content)
                    text_content.append("")  # Add spacing between slides

            return {
                'content': '\n'.join(text_content),
                'metadata': {
                    'slides': len(prs.slides)
                }
            }
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {str(e)}")
            return {'error': str(e)}

    def _process_text(self, file_path: str, mime_type: str) -> Dict[str, Any]:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            return {
                'content': content,
                'metadata': {
                    'size': os.path.getsize(file_path),
                    'lines': len(content.splitlines())
                }
            }
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {str(e)}")
            return {'error': str(e)}

    def _process_code(self, file_path: str) -> Dict[str, Any]:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            return {
                'content': content,
                'metadata': {
                    'language': os.path.splitext(file_path)[1][1:],
                    'lines': len(content.splitlines())
                }
            }
        except Exception as e:
            logger.error(f"Error processing code file {file_path}: {str(e)}")
            return {'error': str(e)}